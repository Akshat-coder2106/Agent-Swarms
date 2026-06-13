from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

def add_slide(title, body):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    body_placeholder.text = body
    
    # Adjust font sizes for readability
    for paragraph in body_placeholder.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        
    return slide

# SLIDE 1
slide_layout_title = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide_layout_title)
slide1.shapes.title.text = "Project Sentinel\nAutonomous DevSecOps Swarm"
slide1.placeholders[1].text = "From vulnerability alert to validated, human-approved fix — automatically.\n\nBuilt on: Azure AI Foundry · Semantic Kernel · AutoGen · Azure OpenAI · GitHub Actions\n\nLive Dashboard: https://agent-swarms-2.vercel.app"

# SLIDE 2
add_slide("Finding vulnerabilities is solved.\nFixing them safely is not.",
    "Alert Overload\n"
    "SAST and DAST tools generate thousands of findings per week. Security teams triage, not fix.\n\n"
    "Remediation Bottleneck\n"
    "Every alert requires a human engineer to analyse the code, write a fix, test it, and submit a PR. That loop takes days, not minutes.\n\n"
    "Business Impact\n"
    "Blocked CI/CD pipelines. Delayed releases. Alert fatigue where critical issues get ignored alongside low-priority noise."
)

# SLIDE 3
add_slide("Sentinel closes the loop.",
    "Sentinel ingests your repository, generates a context-aware patch using Azure OpenAI, validates it inside a process-isolated sandbox, and surfaces one decision to a human engineer: Approve or Rollback.\n\n"
    "Four agents work in sequence:\n"
    "— Architect maps the repository data flow before any patch is written\n"
    "— Scout retrieves semantic context and CVE history for the finding\n"
    "— Engineer generates the patch with full code context\n"
    "— Critic compiles, runs tests, and adversarially challenges the fix in the sandbox\n\n"
    "Nothing reaches a human until the sandbox approves it."
)

# SLIDE 4
add_slide("Built natively on the Microsoft AI stack.",
    "What we use:\n"
    "· Azure OpenAI (gpt-4o) — patch generation and Critic adversarial analysis\n"
    "· Azure AI Foundry — semantic code search via text-embedding-3-small\n"
    "· Semantic Kernel — 4 agent functions registered as @kernel_function plugin\n"
    "· AutoGen GroupChat — 4-agent audit conversation transcript per session\n"
    "· GitHub Actions — one-file CI integration via action.yml\n"
    "· SARIF 2.1 — findings export to GitHub Advanced Security\n"
    "· Azure DevOps — findings exported as work items via REST API\n\n"
    "Live endpoints judges can call:\n"
    "GET  /api/system/sk_status        → Semantic Kernel kernel state\n"
    "GET  /api/sessions/{id}/autogen_transcript → AutoGen conversation log\n"
    "GET  /api/sessions/{id}/export/sarif       → SARIF 2.1 output\n"
    "POST /api/sessions/{id}/export/azure       → Azure DevOps work items"
)

# SLIDE 5
add_slide("Four agents. One convergence loop.\nZero unvalidated patches in production.",
    "The LangGraph state machine drives deterministic retry logic — if the Critic rejects a patch, the Engineer refactors and tries again, up to 7 iterations. Azure AI Foundry embeddings give every agent semantic context, not just keyword matching. Semantic Kernel exposes each agent as a callable plugin function.\n\n"
    "Frontend (Vercel) → JWT Auth → FastAPI Backend (Render)\n"
    "                               ↓\n"
    "                    Semantic Kernel Kernel\n"
    "                    [Architect | Scout | Engineer | Critic]\n"
    "                               ↓\n"
    "                    LangGraph State Machine (retry loop)\n"
    "                               ↓\n"
    "                    Process-Isolated Sandbox\n"
    "                               ↓\n"
    "                    Human Approval Gate → GitHub PR"
)

# SLIDE 6
add_slide("The agents don't just act — they reason, debate, and sign off.",
    "Every audit session produces an AutoGen GroupChat transcript — a readable record of how each agent assessed the vulnerability and the proposed fix.\n\n"
    "Architect: 'Repository analysis complete. SQL injection in users.py line 42 via f-string query construction. Data flow confirmed from untrusted HTTP parameter to database cursor.'\n\n"
    "Scout: 'CWE-89 confirmed. Deterministic rule python.sql_injection.fstring matched. No encoding bypass possible via this code path — parameterisation required.'\n\n"
    "Engineer: 'Patch applied: replaced f-string query with parameterised cursor.execute(). Confidence 0.91. Minimal footprint — one line changed.'\n\n"
    "Critic: 'Sandbox result: APPROVE. Tests pass. Encoding bypass and second-order injection vectors tested — both blocked by parameterisation. Human review recommended before merge.'"
)

# SLIDE 7
add_slide("Tested against 5 real vulnerable repositories.\n86% sandbox pass rate. 223 findings. 150 patches generated.",
    "Repository | Findings | Patches | Pass Rate | False Positives\n"
    "examples/python-vulnerable-api | 14 | 14 | 100% | 0\n"
    "PyCQA/bandit test fixtures | 48 | 42 | 87% | 3\n"
    "trufflesecurity/trufflehog | 32 | 31 | 96% | 1\n"
    "juice-shop/juice-shop | 105 | 45 | 68% | 12\n"
    "appsecco/dvna | 24 | 18 | 80% | 2\n"
    "Overall | 223 | 150 | ~86% | 18\n\n"
    "Note:\n"
    "18 of 223 findings were flagged as false positives and escalated for human review — consistent with our policy gate design. The sandbox rejected 15% of LLM-generated patches before they could reach production."
)

# SLIDE 8
add_slide("Autonomous patching with hard guardrails.\nNo patch merges without passing three policy checks.",
    "Sentinel's policy gate enforces three conditions before any patch is eligible for auto-approval:\n\n"
    "1. Sandbox verdict must be APPROVE — the patch compiled and all tests passed\n"
    "2. Engineer confidence must be ≥ 0.92 — below this threshold, human review is mandatory\n"
    "3. Risk level must be below CRITICAL or HIGH — high-risk patches always require a human\n\n"
    "If any condition fails, the patch is escalated — never silently dropped or auto-merged.\n\n"
    "Rollback is always available: Sentinel stores the original file content and verifies it before restoring, so engineers can revert with one click even after approval."
)

# SLIDE 9
add_slide("11 deterministic detection rules.\nCWE-mapped. OWASP Top 10 tagged. NIST SP 800-53 aligned.",
    "SQL injection via f-string | Python | CWE-89 | A03:2021 | HIGH\n"
    "Code injection via eval()/exec() | Python | CWE-94 | A03:2021 | CRITICAL\n"
    "Path traversal | Python | CWE-22 | A01:2021 | HIGH\n"
    "Insecure deserialization (pickle) | Python | CWE-502 | A08:2021 | CRITICAL\n"
    "Unsafe YAML load | Python | CWE-502 | A08:2021 | HIGH\n"
    "Weak PRNG for security | Python | CWE-338 | A02:2021 | MEDIUM\n"
    "SQL injection via template literal | JavaScript | CWE-89 | A03:2021 | HIGH\n"
    "XSS via innerHTML | JavaScript | CWE-79 | A03:2021 | HIGH\n"
    "Code injection via eval() | JavaScript | CWE-94 | A03:2021 | CRITICAL\n"
    "Hardcoded credentials | Multi | CWE-798 | A07:2021 | CRITICAL\n"
    "Git merge conflict markers | Multi | CWE-116 | A05:2021 | LOW\n\n"
    "Plus: optional Semgrep, Trivy, Checkov, and Gitleaks integration when installed."
)

# SLIDE 10
add_slide("The foundation is in place.\nThree phases to production scale.",
    "Phase 1 — Fleet Scanning\n"
    "Extend the session model to scan multiple repositories in one operation. Useful for organisations auditing 50+ microservices.\n\n"
    "Phase 2 — Expanded Rule Coverage\n"
    "Broader deterministic rules covering complex injection chains, auth bypass patterns, and infrastructure-as-code misconfigurations via Checkov integration.\n\n"
    "Phase 3 — Natural Language Policy Engine\n"
    "Let security teams define policies in plain English: 'Always escalate CRITICAL findings in payment services.' No YAML. No DSL. Powered by Azure OpenAI.\n\n"
    "CLOSING STATEMENT:\n"
    "Sentinel turns security alerts into reviewed, validated, rollback-safe fixes — with full human control at every gate.\n\n"
    "Live demo: https://agent-swarms-2.vercel.app\n"
    "GitHub: https://github.com/Akshat-coder2106/Agent-Swarms\n"
    "API health: https://agent-swarms-api.onrender.com/api/health"
)

prs.save("Project_Sentinel_Presentation.pptx")
